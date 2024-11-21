from pptx import Presentation
from pptx.util import Inches, Pt
import os
import logging

class PowerPointGenerator:
    def __init__(self, ruta_base, img_width_cm=20.32, img_height_cm=9, left_cm=2.54, top_cm=[0.35, 9.53], font_size=72):
        logging.info('Iniciando PowerPointGenerator...')
        self.ruta_base = ruta_base
        self.img_width = Inches(img_width_cm / 2.54)  # Convertir a pulgadas
        self.img_height = Inches(img_height_cm / 2.54)  # Convertir a pulgadas
        self.left_position = Inches(left_cm / 2.54)  # Convertir a pulgadas
        self.top_positions = [Inches(pos / 2.54) for pos in top_cm]  # Convertir a pulgadas
        self.font_size = Pt(font_size)  # Tamaño de la fuente
        self.ppt = Presentation()  # Crear la presentación
        self.imagenes_por_slide = [
            ['ExtAcum_CuT_%_vs_dias', 'Ext%Cu_IL_vs_dias'],
            ['Fe+3_FeT_PLS_vs_dias', 'ORP_PLS_mV_vs_dias'],
            ['Ext%Cu_IL_vs_acido_acum_agregado_refino_kg_t', 'Ext%Cu_IL_vs_ConsNeto_Aci_kg_t'],
            ['pH_PLS_vs_dias', 'acido_PLS_g_kg_vs_dias']
        ]

    def agregar_slide_con_titulo(self, grupo):
        """Crea una diapositiva de título para cada grupo."""
        try:
            logging.info(f'Agregando diapositiva de título para el grupo {grupo}...')
            slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[5])
            title = slide.shapes.title
            title.text = f'Grupo {grupo}'
            title.text_frame.paragraphs[0].font.size = self.font_size
            title.text_frame.paragraphs[0].alignment = 1  # Centrar el texto
        except Exception as e:
            logging.error(f'Error al agregar diapositiva de título para el grupo {grupo}: {e}')

    def agregar_imagenes(self, carpeta_actual):
        """Añade imágenes a la diapositiva según las posiciones definidas."""
        try:
            logging.info(f'Agregando imágenes desde {carpeta_actual}...')
            for grupo_imagenes in self.imagenes_por_slide:
                slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[6])  # Layout en blanco
                for img_name, top in zip(grupo_imagenes, self.top_positions):
                    imagen_encontrada = False
                    for imagen in sorted(os.listdir(carpeta_actual)):
                        if img_name in imagen and (imagen.endswith('.png') or imagen.endswith('.jpg') or imagen.endswith('.jpeg')):
                            img_path = os.path.join(carpeta_actual, imagen)
                            slide.shapes.add_picture(img_path, self.left_position, top, width=self.img_width, height=self.img_height)
                            imagen_encontrada = True
                            logging.info(f'Imagen {imagen} agregada a la diapositiva desde {img_path}')
                            break  # Salir del bucle al encontrar la imagen correspondiente
                    if not imagen_encontrada:
                        logging.warning(f'Imagen con nombre {img_name} no encontrada en {carpeta_actual}')
        except Exception as e:
            logging.error(f'Error al agregar imágenes desde {carpeta_actual}: {e}')

    def generar_presentacion(self, total_grupos=12, output_name='Presentación_columnas.pptx'):
        """Genera la presentación completa con las imágenes para todos los grupos."""
        try:
            logging.info('Iniciando generación de la presentación...')
            for i in range(1, total_grupos + 1):
                carpeta_actual = os.path.join(self.ruta_base, f'graficas_grupo{i}')
                if os.path.exists(carpeta_actual):  # Verificar que la carpeta exista
                    self.agregar_slide_con_titulo(i)
                    self.agregar_imagenes(carpeta_actual)
                else:
                    logging.warning(f'Carpeta no encontrada: {carpeta_actual}')

            # Guardar la presentación
            output_path = os.path.join(self.ruta_base, output_name)
            self.ppt.save(output_path)
            logging.info(f'Presentación guardada en {output_path}')
            print(f'Presentación guardada en {output_path}')
        except Exception as e:
            logging.error(f'Error al generar la presentación: {e}')
            print(f'Error al generar la presentación: {e}')

